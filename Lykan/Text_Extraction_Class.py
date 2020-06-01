# -*- coding: utf-8 -*-
"""
Created on Mon Jun  1 15:53:40 2020

@author: Devil~Gareeeb7
"""   
import glob
import cv2
import PyPDF2
import pytesseract
from pptx import Presentation
from docx import Document
import pandas


if __name__ == "__main__":

    class Text_Extraction:

# pdf file text       
        def pdf_To_Text(file = " "):     # pdf file text

            pdfFileObj = open(file, 'rb') 
            pdfReader = PyPDF2.PdfFileReader(pdfFileObj)  
            pageObj = pdfReader.getPage(0) 
            pdf_file = pageObj.extractText()
        
            return pdf_file


# docx file text
        def docx_To_Text(file = " "):
        
            document = Document(file)
            for para in document.paragraphs:
                docx_file = para.text
                docx_file
        
            return docx_file


# text file text
        def load_Text(file = " "):
        
            myfile = open(file, "rt") 
            text = myfile.read()
        
            return text


# ppt file text
        def ppt_To_Text(file = " "):
        
            for eachfile in glob.glob(file):
                prs = Presentation(eachfile)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                                ppt_file = shape.text

            return ppt_file


# image file text
        def image_To_Text(file = " "):

            img = cv2.imread(file)
            img2 = cv2.resize(img,None,fx = 2, fy = 2)

            text = pytesseract.image_to_string(img2, lang = 'eng')
        
            return text

# csv data loading
        
        def csv_Load(file = " "):
        
            data = pandas.read_csv(file)
        
            return data




